import hashlib
import os
import shutil
import uuid
from datetime import datetime, timezone

from PyPDF2 import PdfReader
from PyPDF2.errors import PdfReadError
from pptx import Presentation
from docx import Document
from sqlalchemy import delete, select

from app.config import settings
from app.database import async_session
from app.models.uploaded_file import UploadedFile
from app.models.document_chunk import DocumentChunk
from app.models.enums import FileStatus
from app.services.chroma_store import add_chunks_to_chroma, delete_document_vectors

ALLOWED_EXTENSIONS = {".pdf", ".docx", ".ppt", ".pptx"}
MIME_MAP = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".ppt": "application/vnd.ms-powerpoint",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

TMP_DIR = os.path.join(settings.UPLOAD_DIR, "tmp")


def save_uploaded_file(file_content: bytes, original_filename: str) -> tuple[str, str, str]:
    ext = os.path.splitext(original_filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        ext = ".bin"
    file_id = str(uuid.uuid4())
    tmp_filename = f"{file_id}{ext}"
    os.makedirs(TMP_DIR, exist_ok=True)
    tmp_path = os.path.join(TMP_DIR, tmp_filename)

    sha256 = hashlib.sha256()
    with open(tmp_path, "wb") as f:
        sha256.update(file_content)
        f.write(file_content)

    return tmp_path, sha256.hexdigest(), ext


def move_to_category_path(tmp_path: str, category: str, user_id: str, doc_id: str) -> str:
    ext = os.path.splitext(tmp_path)[1]
    target_dir = os.path.join(settings.UPLOAD_DIR, category, user_id)
    os.makedirs(target_dir, exist_ok=True)
    target_path = os.path.join(target_dir, f"{doc_id}{ext}")
    shutil.move(tmp_path, target_path)
    return target_path


class EmptyDocumentError(Exception):
    pass


class EncryptedDocumentError(Exception):
    pass


class UnsupportedFormatError(Exception):
    pass


MIN_TEXT_CHARS = 50


def _extract_pdf_text(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
    except PdfReadError:
        raise EncryptedDocumentError("该 PDF 已加密，无法提取文本")

    if reader.is_encrypted:
        raise EncryptedDocumentError("该 PDF 已加密，无法提取文本")

    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)

    full_text = "\n".join(parts).strip()
    if not full_text:
        raise EmptyDocumentError("该文档无可提取文本，请上传文字版 PDF")
    return full_text


def _extract_docx_text(file_path: str) -> str:
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    full_text = "\n".join(parts).strip()
    if not full_text:
        raise EmptyDocumentError("该文档无可提取文本，请上传文字版 DOCX")
    return full_text


def _extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
    full_text = "\n".join(parts).strip()
    if not full_text:
        raise EmptyDocumentError("该文档无可提取文本，请上传文字版 PPTX")
    return full_text


def extract_text(file_path: str, mime_type: str) -> str:
    if mime_type == "application/pdf":
        return _extract_pdf_text(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return _extract_docx_text(file_path)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_pptx_text(file_path)
    elif mime_type == "application/vnd.ms-powerpoint":
        raise UnsupportedFormatError("不支持旧版 PPT 格式，请转换为 PPTX 后上传")
    else:
        raise UnsupportedFormatError(f"不支持的文件格式: {mime_type}")


async def process_file_extraction(file_id: str):
    async with async_session() as session:
        result = await session.execute(
            select(UploadedFile).where(UploadedFile.id == file_id)
        )
        file = result.scalar_one_or_none()
        if not file:
            return

        try:
            text = extract_text(file.path, file.mimeType)

            if len(text) < MIN_TEXT_CHARS:
                file.status = FileStatus.failed
                file.errorMessage = "提取文字过少，可能为图片文档"
                await session.commit()
                return

            file.parsedText = text
            file.totalChars = len(text)
            file.parsedAt = datetime.utcnow()
            file.status = FileStatus.parsed
            await session.commit()

        except EmptyDocumentError as e:
            file.status = FileStatus.failed
            file.errorMessage = str(e)
            await session.commit()
            return

        except EncryptedDocumentError as e:
            file.status = FileStatus.failed
            file.errorMessage = str(e)
            await session.commit()
            return

        except UnsupportedFormatError as e:
            file.status = FileStatus.failed
            file.errorMessage = str(e)
            await session.commit()
            return

        except Exception as e:
            file = await session.get(UploadedFile, file_id)
            if file:
                file.status = FileStatus.failed
                file.errorMessage = str(e)
                await session.commit()
            return

    await _classify_file(file_id)

    async with async_session() as session:
        file = await session.get(UploadedFile, file_id)
        if not file or file.status == FileStatus.failed:
            return

    await _chunk_file(file_id)

    async with async_session() as session:
        file = await session.get(UploadedFile, file_id)
        if not file or file.status == FileStatus.failed:
            return

    async with async_session() as session:
        file = await session.get(UploadedFile, file_id)
        if file:
            await process_file_embedding(file_id, file.category or "general", file.userId)


async def _classify_file(file_id: str):
    from app.services.classifier import classify_and_move_file

    async with async_session() as session:
        file = await session.get(UploadedFile, file_id)
        if not file or not file.parsedText:
            return

        try:
            await classify_and_move_file(session, file)
        except Exception:
            file = await session.get(UploadedFile, file_id)
            if file:
                file.category = "general"
                file.categoryConfidence = 0.0
                await session.commit()


async def _chunk_file(file_id: str):
    from app.services.chunker import chunk_document

    async with async_session() as session:
        file = await session.get(UploadedFile, file_id)
        if not file or not file.parsedText:
            return

        try:
            await chunk_document(session, file_id, file.parsedText)
        except Exception:
            file = await session.get(UploadedFile, file_id)
            if file:
                file.status = FileStatus.failed
                file.errorMessage = "chunking failed"
                await session.commit()




async def delete_file_record(file_id: str, user_id: str) -> dict | None:
    async with async_session() as session:
        result = await session.execute(
            select(UploadedFile).where(
                UploadedFile.id == file_id,
                UploadedFile.userId == user_id,
            )
        )
        file = result.scalar_one_or_none()
        if not file:
            return None

        if file.category and file.userId:
            try:
                await delete_document_vectors(file_id, file.category, file.userId)
            except Exception:
                pass

        if file.path and os.path.exists(file.path):
            try:
                os.remove(file.path)
            except Exception:
                pass

        original_name = file.originalName

        await session.execute(
            delete(UploadedFile).where(UploadedFile.id == file_id)
        )
        await session.commit()

        return {"id": file_id, "original_name": original_name}


async def process_file_embedding(file_id: str, category: str, user_id: str):
    async with async_session() as session:
        result = await session.execute(
            select(UploadedFile).where(UploadedFile.id == file_id)
        )
        file = result.scalar_one_or_none()
        if not file:
            return

        try:
            file.status = FileStatus.embedding
            await session.commit()

            result = await session.execute(
                select(DocumentChunk)
                .where(DocumentChunk.fileId == file_id)
                .order_by(DocumentChunk.chunkIndex)
            )
            chunks = result.scalars().all()

            if not chunks:
                file.status = FileStatus.ready
                file.chunkCount = 0
                await session.commit()
                return

            chunk_dicts = [
                {
                    "content": c.content,
                    "fileId": c.fileId,
                    "chunkIndex": c.chunkIndex,
                }
                for c in chunks
            ]

            await add_chunks_to_chroma(
                category=category,
                user_id=user_id,
                chunks=chunk_dicts,
            )

            file.status = FileStatus.ready
            file.chunkCount = len(chunks)
            await session.commit()

        except Exception as e:
            file = await session.get(UploadedFile, file_id)
            if file:
                file.status = FileStatus.failed
                file.errorMessage = str(e)
                await session.commit()
